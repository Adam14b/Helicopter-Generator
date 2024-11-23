from report import Report
from pptx.enum.shapes import MSO_SHAPE
from datetime import date
from config import *
from copy import deepcopy

class Timeline:
    def __init__(self, report: Report, y_top: int, y_bottom: int, start_date: date, final_date: date, description: list = [], show_start_date=True) -> None:
        self.report = report
        self.timeline_y = (y_top + y_bottom) / 2
        self.y_top = y_top
        self.y_bottom = y_bottom

        try:
            start_x = report.date_to_x(start_date)
            final_x = report.date_to_x(final_date)
        except ValueError:
            return

        try:
            now_x = report.date_to_x(report.now_date)
        except ValueError:
            now_x = None

        if final_date < report.now_date:
            report.add_shape(MSO_SHAPE.OVAL, start_x - 6, self.timeline_y - 6, 12, 12, color=GREEN, fill=True)
            report.add_shape(MSO_SHAPE.LINE_INVERSE, start_x + 6, self.timeline_y, final_x - start_x - 6, 0, color=GREEN)
        elif start_date < report.now_date <= final_date:
            report.add_shape(MSO_SHAPE.OVAL, start_x - 6, self.timeline_y - 6, 12, 12, color=GREEN, fill=True)
            if now_x:
                report.add_shape(MSO_SHAPE.LINE_INVERSE, start_x + 6, self.timeline_y, now_x - start_x - 6, 0, color=GREEN)
                report.add_shape(MSO_SHAPE.LINE_INVERSE, now_x, self.timeline_y, final_x - now_x, 0, color=GRAY)
            else:
                report.add_shape(MSO_SHAPE.LINE_INVERSE, start_x + 6, self.timeline_y, final_x - start_x - 6, 0, color=GRAY)
        elif report.now_date <= start_date:
            report.add_shape(MSO_SHAPE.OVAL, start_x - 6, self.timeline_y - 6, 12, 12, color=GRAY, fill=True, fill_color=WHITE)
            report.add_shape(MSO_SHAPE.LINE_INVERSE, start_x + 6, self.timeline_y, final_x - start_x - 6, 0, color=GRAY)

        if show_start_date:
            report.add_text(date_to_text(start_date), start_x - 16, self.timeline_y + 10, 32, 10)
        report.add_text(date_to_text(final_date), final_x - 16, self.timeline_y + 10, 32, 10)

        report.add_shape(MSO_SHAPE.LINE_INVERSE, 30, y_top, 930, 0, color=LIGHTBLUE, fill=True)
        report.add_task(description, 30, y_top, 210, abs(y_bottom - y_top))

    def add_figure(self, figure_description: str, date_: date, write_date: bool = False, note: str = ''):
        if '/' in figure_description:
            figure_type, color_str = tuple(figure_description.split('/'))
            color = color_dict.get(color_str, GREEN)
        else:
            figure_type = figure_description
            color = None

        try:
            x_position = self.report.date_to_x(date_)
        except ValueError:
            return

        if figure_type in ['план', 'сдвиг срока', 'провал']:
            if color is None:
                if figure_type == 'план' and date_ <= self.report.now_date:
                    color = GREEN
                elif figure_type == 'план' and date_ > self.report.now_date:
                    color = GRAY
                elif figure_type == 'сдвиг срока':
                    color = YELLOW
                elif figure_type == 'провал':
                    color = RED
            self.report.add_figure(colored_flag(color, date_ <= self.report.now_date), x_position - 6, self.timeline_y - 12)
        else:
            figure = deepcopy(fig_dict[figure_type])
            if color is None:
                color = GREEN if date_ <= self.report.now_date else GRAY
            if isinstance(figure, Shape):
                figure.fill.solid()
                figure.fill.fore_color.rgb = color
            elif isinstance(figure, GroupShape):
                for subfigure in figure.shapes:
                    if isinstance(subfigure, Shape):
                        subfigure.fill.solid()
                        subfigure.fill.fore_color.rgb = color

            self.report.add_figure(figure, x_position - 6, self.timeline_y - 14)

        if write_date:
            self.report.add_text(date_to_text(date_), x_position - 20, self.timeline_y - 20, width_pt=40, height_pt=10)
        if note:
            self.report.add_text(note, x_position - 20, self.timeline_y + 3, width_pt=40, height_pt=10)

    def add_arrow(self, start_date: date, final_date: date):
        try:
            start_x = self.report.date_to_x(start_date)
            final_x = self.report.date_to_x(final_date)
        except ValueError:
            return
        self.report.add_arrow(start_x - 6, self.timeline_y - 3, final_x - start_x, 0, color=YELLOW if start_date < final_date else GREEN)

    def add_gold(self, text: str, y: int):
        shape = self.report.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, 0, self.y_top + y, 30, 10, fill=True, color=BLACK, fill_color=SUNNY_YELLOW)
        shape.line.width = Pt(0)
        shape_paragraph = shape.text_frame.paragraphs[0]
        shape_paragraph.text = text
        shape_paragraph.font.size = Pt(5)
        shape_paragraph.font.name = 'Arial'
        shape_paragraph.font.color.rgb = BLACK

    def add_comment(self, text: str, date_: date, size: float):
        try:
            x_position = self.report.date_to_x(date_)
        except ValueError:
            return
        shape = self.report.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x_position - size / 2, self.timeline_y + 10, size, 10, fill=False, color=YELLOW)
        shape_paragraph = shape.text_frame.paragraphs[0]
        shape_paragraph.text = text
        shape_paragraph.font.size = Pt(5)
        shape_paragraph.font.name = 'Arial'
        shape_paragraph.font.color.rgb = BLUE
