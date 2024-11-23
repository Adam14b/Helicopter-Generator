from datetime import date, timedelta
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.dml import MSO_LINE 
from pptx.oxml import parse_xml
from pptx import Presentation
from config import *
import numpy as np
import os

class Report:
    def __init__(self, now_date: date, quarters_with_tasks, tasks_per_quarter, total_tasks, min_quarter_width=60, max_quarter_width=300, left_margin=240, right_margin=20) -> None:
        self.presentation = Presentation()
        self.presentation.slide_width = Pt(960)
        self.presentation.slide_height = Pt(540)
        self.slide = self.presentation.slides.add_slide(self.presentation.slide_layouts[6])

        self.now_date = now_date
        self.quarters_with_tasks = quarters_with_tasks
        self.tasks_per_quarter = tasks_per_quarter
        self.total_tasks = total_tasks

        self.quarter_list = quarters_with_tasks
        available_width = 960 - left_margin - right_margin
        total_weight = sum(tasks_per_quarter.values())
        self.kvartal_pixels = []
        for quarter in self.quarter_list:
            weight = tasks_per_quarter[quarter] / total_weight
            width = available_width * weight
            width = max(min_quarter_width, min(width, max_quarter_width))
            self.kvartal_pixels.append(width)

        total_quarter_width = sum(self.kvartal_pixels)
        if total_quarter_width > available_width:
            scaling_factor = available_width / total_quarter_width
            self.kvartal_pixels = [width * scaling_factor for width in self.kvartal_pixels]
        else:
            extra_space = available_width - total_quarter_width
            self.kvartal_pixels = [width + extra_space / len(self.kvartal_pixels) for width in self.kvartal_pixels]

        self.kvartal_pixel_bounds = np.cumsum([0] + self.kvartal_pixels)
        self.kvartal_days = []
        for year, quarter in self.quarter_list:
            q_start_month = (quarter - 1) * 3 + 1
            q_end_month = q_start_month + 2
            q_start = date(year, q_start_month, 1)
            if q_end_month >= 12:
                q_end = date(year, 12, 31)
            else:
                q_end = date(year, q_end_month + 1, 1) - timedelta(days=1)
            days_in_quarter = (q_end - q_start).days + 1
            self.kvartal_days.append(days_in_quarter)
        self.kvartal_day_bounds = np.cumsum([0] + self.kvartal_days)

        self.start_year = self.quarter_list[0][0]
        self.start_quarter = self.quarter_list[0][1]
        self.left_margin = left_margin
        self.create_base_structure()

    def create_base_structure(self):
        self.add_shape(MSO_SHAPE.RECTANGLE, 0, 60, 960, 30, color=BLUE)
        for i in range(len(self.kvartal_pixel_bounds)):
            x_position = self.left_margin + self.kvartal_pixel_bounds[i]
            self.add_shape(MSO_SHAPE.LINE_INVERSE, x_position, 60, 1, 480, color=LIGHTBLUE)

        self.add_shape(MSO_SHAPE.LINE_INVERSE, 15, 60, 1, 480, color=LIGHTBLUE)
        self.add_shape(MSO_SHAPE.LINE_INVERSE, 30, 60, 1, 480, color=LIGHTBLUE)

        logo_path = os.path.join(CONTENT_DIR, 'vtb_logo.png')
        self.slide.shapes.add_picture(logo_path, Pt(880), Pt(10), Pt(127 // 2), Pt(52 // 2))

        self.add_text(f'Кластер «Управление Продажами»', 20, 0, size_pt=20, bold=True, width_pt=940, align=PP_ALIGN.LEFT)

        self.add_text('№', 0, 60, 15, 30, color=WHITE, size_pt=8, bold=True, anchor=MSO_ANCHOR.MIDDLE)
        self.add_text('К', 15, 60, 15, 30, color=WHITE, size_pt=8, bold=True, anchor=MSO_ANCHOR.MIDDLE)
        self.add_text('Задача', 30, 60, 210, 30, color=WHITE, size_pt=8, bold=True, anchor=MSO_ANCHOR.MIDDLE)

        months = ['январь', 'февраль', 'март', 'апрель', 'май', 'июнь', 'июль', 'август', 'сентябрь', 'октябрь', 'ноябрь', 'декабрь']
        for i, (year, quarter) in enumerate(self.quarter_list):
            quarter_name = f'{quarter}Q {year}'
            x_position = self.left_margin + self.kvartal_pixel_bounds[i]
            width = self.kvartal_pixels[i]
            self.add_text(quarter_name, x_position + 1, 60, width, 15, color=WHITE, size_pt=8, bold=True, anchor=MSO_ANCHOR.MIDDLE)
            for m in range(3):
                month_index = (quarter - 1) * 3 + m
                if month_index >= 12:
                    break
                month_name = months[month_index]
                month_start_date = date(year, month_index + 1, 1)
                x_position_month = self.date_to_x(month_start_date)
                self.add_text(month_name, x_position_month, 75, self.kvartal_pixels[i]/3, 15, color=WHITE, size_pt=6, bold=True, anchor=MSO_ANCHOR.MIDDLE)

        try:
            now_x = self.date_to_x(self.now_date)
            now_line = self.add_shape(MSO_SHAPE.LINE_INVERSE, now_x, 90, 1, 410, color=RED, fill=True)
            now_line.line.width = Pt(2)
            now_line.line.dash_style = MSO_LINE.DASH
            self.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, now_x - 2.7, 500, 5.4, 8, color=BLUE, fill=True)
            self.add_text(date_to_text(self.now_date), now_x - 20, 510, color=RED, size_pt=9, width_pt=40, height_pt=20, bold=True)
        except ValueError:
            pass

        self.add_legend()

    def add_legend(self):
        self.add_figure(key_fig, 0, 520)
        self.add_text('ИФТ', 15, 522, width_pt=120, height_pt=16, size_pt=8, color=BLACK, align=PP_ALIGN.LEFT)
        self.add_figure(gear_fig, 40, 520)
        self.add_text('НТ', 55, 522, width_pt=120, height_pt=16, size_pt=8, color=BLACK, align=PP_ALIGN.LEFT)
        self.add_figure(box_fig, 80, 520)
        self.add_text('ПСИ', 95, 522, width_pt=120, height_pt=16, size_pt=8, color=BLACK, align=PP_ALIGN.LEFT)
        self.add_figure(rocket_fig, 120, 520)
        self.add_text('Прод', 135, 522, width_pt=120, height_pt=16, size_pt=8, color=BLACK, align=PP_ALIGN.LEFT)
        self.add_figure(pin_fig, 160, 520)
        self.add_text('MVP', 175, 522, width_pt=120, height_pt=16, size_pt=8, color=BLACK, align=PP_ALIGN.LEFT)
        self.add_figure(pilot_fig, 200, 520)
        self.add_text('Пилот', 215, 522, width_pt=120, height_pt=16, size_pt=8, color=BLACK, align=PP_ALIGN.LEFT)
        self.add_figure(colored_flag(GRAY, filled=False), 490, 520)
        self.add_text('План', 500, 522, width_pt=200, height_pt=16, size_pt=8, color=BLACK, align=PP_ALIGN.LEFT)
        self.add_arrow(530, 525, 30, 0)
        self.add_text('Перенос срока', 565, 522, width_pt=200, height_pt=16, size_pt=8, color=BLACK, align=PP_ALIGN.LEFT)
        self.add_figure(colored_flag(GREEN, filled=True), 640, 520)
        self.add_text('Выполнено', 650, 522, width_pt=200, height_pt=16, size_pt=8, color=BLACK, align=PP_ALIGN.LEFT)
        self.add_figure(colored_flag(YELLOW, filled=False), 700, 520)
        self.add_text('Риск сдвига', 710, 522, width_pt=200, height_pt=16, size_pt=8, color=BLACK, align=PP_ALIGN.LEFT)
        self.add_figure(colored_flag(RED, filled=False), 760, 520)
        self.add_text('Просрочено', 770, 522, width_pt=200, height_pt=16, size_pt=8, color=BLACK, align=PP_ALIGN.LEFT)
        self.add_figure(colored_flag(YELLOW, filled=True), 820, 520)
        self.add_text('Выполнено со сдвигом срока', 830, 522, width_pt=200, height_pt=16, size_pt=8, color=BLACK, align=PP_ALIGN.LEFT)

    def add_shape(self, shape: MSO_SHAPE_TYPE, left_pt, top_pt, width_pt, height_pt, color=GREEN, fill=True, fill_color=None):
        shape_obj = self.slide.shapes.add_shape(shape, Pt(left_pt), Pt(top_pt), Pt(width_pt), Pt(height_pt))
        if fill:
            shape_obj.fill.solid()
            shape_obj.fill.fore_color.rgb = color if fill_color is None else fill_color
        else:
            shape_obj.fill.background()
            shape_obj.line.width = Pt(1)
        if fill_color is not None:
            shape_obj.line.width = Pt(2)
        shape_obj.line.color.rgb = color
        return shape_obj

    def add_arrow(self, left_pt, top_pt, width_pt, height_pt, color=YELLOW):
        shape = self.slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Pt(left_pt + width_pt), Pt(top_pt), Pt(left_pt), Pt(top_pt + height_pt))
        shape.line._get_or_add_ln().append(parse_xml('<a:headEnd type="arrow" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"/>'))
        shape.line.width = Pt(2)
        shape.line.dash_style = MSO_LINE.DASH
        shape.line.fill.solid()
        shape.line.color.rgb = color
        return shape

    def save(self, filename: str):
        self.presentation.save(filename)

    def date_to_x(self, date_: date):
        try:
            kv = self.quarter_list.index((date_.year, (date_.month - 1) // 3 + 1))
        except ValueError:
            raise ValueError("Date out of range")
        kvartal_pixel = self.kvartal_pixels[kv]
        kvartal_day = self.kvartal_days[kv]
        q_year, q_quarter = self.quarter_list[kv]
        q_start_month = (q_quarter - 1) * 3 + 1
        q_start_date = date(q_year, q_start_month, 1)
        days_into_quarter = (date_ - q_start_date).days
        x = self.left_margin + self.kvartal_pixel_bounds[kv] + days_into_quarter * (kvartal_pixel / kvartal_day)
        return x

    def add_figure(self, figure, left_pt, top_pt):
        element = deepcopy(figure.element)
        element.x = Pt(left_pt)
        element.y = Pt(top_pt)
        self.slide.shapes._spTree.append(element)

    def add_text(self, text: str, left_pt, top_pt, width_pt=100, height_pt=40, size_pt=5, color=BLUE, bold=False, anchor=MSO_ANCHOR.TOP, align=PP_ALIGN.CENTER):
        text_frame = self.slide.shapes.add_textbox(Pt(left_pt), Pt(top_pt), Pt(width_pt), Pt(height_pt)).text_frame
        text_frame.margin_left = Pt(0)
        text_frame.margin_top = Pt(0)
        text_frame.margin_right = Pt(0)
        text_frame.margin_bottom = Pt(0)
        text_frame.vertical_anchor = anchor
        textbox_paragraph = text_frame.paragraphs[0]
        textbox_paragraph.text = text
        textbox_paragraph.font.size = Pt(size_pt)
        textbox_paragraph.font.color.rgb = color
        textbox_paragraph.alignment = align
        textbox_paragraph.font.name = 'Arial'
        textbox_paragraph.font.bold = bold
        return textbox_paragraph

    def add_task(self, texts: list, left_pt, top_pt, width_pt=150, height_pt=50):
        texts = [text.replace('//', '\n') for text in texts]
        text_frame = self.slide.shapes.add_textbox(Pt(left_pt), Pt(top_pt), Pt(width_pt - 30), Pt(height_pt)).text_frame
        text_frame.word_wrap = True
        text_frame.margin_left = Pt(3)
        text_frame.margin_top = Pt(1)
        text_frame.margin_right = Pt(1)
        text_frame.margin_bottom = Pt(1)
        text_frame.vertical_anchor = MSO_ANCHOR.TOP
        for idx, text in enumerate(texts[:3]):
            paragraph = text_frame.paragraphs[idx] if idx == 0 else text_frame.add_paragraph()
            paragraph.text = text
            paragraph.font.size = Pt(5)
            paragraph.font.color.rgb = [BLUE, BLUE, DARKGRAY][idx]
            paragraph.alignment = PP_ALIGN.LEFT
            paragraph.font.name = 'Arial'
            paragraph.font.bold = idx == 0
        text_frame_right = self.slide.shapes.add_textbox(Pt(left_pt + width_pt - 30), Pt(top_pt), Pt(30), Pt(height_pt)).text_frame
        text_frame_right.word_wrap = True
        text_frame_right.margin_left = Pt(1)
        text_frame_right.margin_top = Pt(1)
        text_frame_right.margin_right = Pt(1)
        text_frame_right.margin_bottom = Pt(1)
        text_frame_right.vertical_anchor = MSO_ANCHOR.TOP
        for idx, text in enumerate(texts[3:]):
            paragraph = text_frame_right.paragraphs[idx] if idx == 0 else text_frame_right.add_paragraph()
            paragraph.text = text
            paragraph.font.size = Pt(5)
            paragraph.font.color.rgb = [BLUE, YELLOW, GREEN][idx]
            paragraph.alignment = PP_ALIGN.RIGHT
            paragraph.font.name = 'Arial'
