import os
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt
from copy import deepcopy
from pptx.shapes.group import GroupShape
from pptx.shapes.autoshape import Shape

CONTENT_DIR = '../contents'
slide_template_path = os.path.join(CONTENT_DIR, 'examples.pptx')
slide_template = Presentation(slide_template_path).slides[0]

base_flag = deepcopy(slide_template.shapes[0])
green_flag = deepcopy(slide_template.shapes[0])
risky_flag = deepcopy(slide_template.shapes[1])
failed_flag = deepcopy(slide_template.shapes[3])
finally_flag = deepcopy(slide_template.shapes[8])
key_fig = deepcopy(slide_template.shapes[10])
gear_fig = deepcopy(slide_template.shapes[11])
box_fig = deepcopy(slide_template.shapes[12])
rocket_fig = deepcopy(slide_template.shapes[13])
pin_fig = deepcopy(slide_template.shapes[14])
pilot_fig = deepcopy(slide_template.shapes[15])
plan_flag = deepcopy(slide_template.shapes[16])

GREEN = RGBColor(78, 201, 149)
BLUE = RGBColor(10, 40, 150)
LIGHTBLUE = RGBColor(200, 200, 255)
YELLOW = RGBColor(240, 160, 40)
RED = RGBColor(255, 40, 40)
GRAY = RGBColor(194, 194, 194)
WHITE = RGBColor(255, 255, 255)
BLACK = RGBColor(0, 0, 0)
DARKGRAY = RGBColor(106, 106, 106)

base_flag.width = Pt(8.4)
base_flag.height = Pt(12.2)

def colored_flag(color: RGBColor, filled: bool):
    flag: GroupShape = deepcopy(base_flag)
    fabric: Shape = flag.shapes[0]
    stick: Shape = flag.shapes[1]
    fabric.line.color.rgb = color
    fabric.fill.solid()
    fabric.fill.fore_color.rgb = color if filled else WHITE
    stick.line.color.rgb = color
    return flag

def date_to_text(date_):
    return f'{str(date_.day).zfill(2)}.{str(date_.month).zfill(2)}'

def k_to_y(k: int):
    return 90 + 25 + 50*k

def add_text_task(slide, text: str, left_pt, top_pt, width_pt=0, height_pt=0, size_pt=5, rgb_color=BLUE):
    textbox = slide.shapes.add_textbox(Pt(left_pt), Pt(top_pt), Pt(width_pt), Pt(height_pt))
    textbox.text_frame.word_wrap = True
    textbox_paragraph = textbox.text_frame.paragraphs[0]
    textbox_paragraph.text = text
    textbox_paragraph.font.size = Pt(size_pt)
    textbox_paragraph.font.color.rgb = rgb_color
    textbox_paragraph.alignment = PP_ALIGN.LEFT
    textbox_paragraph.font.name = 'Arial'
    textbox_paragraph.font.bold = True
    return textbox_paragraph

FOUNDING = RGBColor(125, 150, 246)
COMPLEXITY = RGBColor(67, 104, 172)

fig_dict = {
    'ИФТ': key_fig,
    'ПСИ': box_fig,
    'НТ': gear_fig,
    'прод': rocket_fig,
    'MVP': pin_fig,
    'пилот': pilot_fig,
    'план': None,
    'сдвиг срока': None,
    'провал': None,
}

color_dict = {
    'зелёный': GREEN,
    'зеленый': GREEN,
    'жёлтый': YELLOW,
    'желтый': YELLOW,
    'красный': RED,
    'серый': GRAY,
}

fig_to_color = {
    'план': GRAY,
    'успех': GREEN,
    'провал': RED,
    'риск переноса': YELLOW,
    'выполнено со сдвигом': YELLOW,
}

SUNNY_YELLOW = RGBColor(255, 255, 0)
