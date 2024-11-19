from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE, MSO_CONNECTOR_TYPE, MSO_CONNECTOR
from pptx.dml.color import RGBColor
from pptx.slide import SlideMasters, Slides, Slide, SlideLayout
from pptx.util import Inches, Pt, Cm, Emu
from copy import deepcopy
from pptx.shapes.group import GroupShape
from pptx.shapes.base import BaseShape
from pptx.shapes.connector import Connector
from pptx.shapes.freeform import FreeformBuilder
from pptx.shapes.autoshape import Shape
from collections import defaultdict
from lxml import etree
import numpy as np
from datetime import datetime, date, timedelta
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.shapes.autoshape import AutoShapeType, Shape, BaseShape
from pptx.enum.dml import MSO_LINE
from pptx.shapes.graphfrm import GraphicFrame
from pptx.text.text import _Paragraph
from pptx.oxml import parse_xml

slide_template = Presentation('examples.pptx').slides[0]  ##

base_flag = deepcopy(slide_template.shapes[0])

green_flag = deepcopy(slide_template.shapes[0])
risky_flag = deepcopy(slide_template.shapes[1])
failed_flag = deepcopy(slide_template.shapes[3])
finally_flag = deepcopy(slide_template.shapes[8])

key_fig = deepcopy(slide_template.shapes[10])
gear_fig = deepcopy(slide_template.shapes[11])
box_fig = deepcopy(slide_template.shapes[12])
rocket_fig = deepcopy(slide_template.shapes[13])
pin_fig = deepcopy(slide_template.shapes[14])
pilot_fig = deepcopy(slide_template.shapes[15])
plan_flag = deepcopy(slide_template.shapes[16])


GREEN = RGBColor(78, 201, 149)
BLUE = RGBColor(10, 40, 150)
LIGHTBLUE = RGBColor(200, 200, 255)
YELLOW = RGBColor(240, 160, 40)
RED = RGBColor(255, 40, 40)
GRAY = RGBColor(194, 194, 194)
WHITE = RGBColor(255, 255, 255)
BLACK = RGBColor(0, 0, 0)
DARKGRAY = RGBColor(106, 106, 106)
SUNNY_YELLOW = RGBColor(255, 255, 0)
FOUNDING = RGBColor(125, 150, 246)
COMPLEXITY = RGBColor(67, 104, 172)

base_flag.width = Pt(8.4)
base_flag.height = Pt(12.2)

def colored_flag(color: RGBColor, filled: bool):
    flag: GroupShape = deepcopy(base_flag)

    fabric: Shape = flag.shapes[0]
    stick: Shape = flag.shapes[1]

    fabric.line.color.rgb = color
    fabric.fill.solid()
    fabric.fill.fore_color.rgb = color if filled else WHITE
    stick.line.color.rgb = color

    return flag


def date_to_x(date_: date, kvartal: int):
    date_tuple = date_.timetuple()
    day_in_year = date_tuple.tm_yday - 1
    month_in_year = date_tuple.tm_mon - 1
    kvartal_pixels = [0] + [60 if k == kvartal else 540 for k in range(4)]
    kvartal_pixel_bounds = np.cumsum(kvartal_pixels)
    kvartal_days = [0] + [90, 91, 92, 92]
    kvartal_day_bounds = np.cumsum(kvartal_days)
    kv = month_in_year // 3

    return kvartal_pixel_bounds[kv] + (day_in_year - kvartal_day_bounds[kv]) * (kvartal_pixel_bounds[kv + 1] / kvartal_days[kv + 1])

def date_to_text(date_: date):
    return f'{str(date_.day).zfill(2)}.{str(date_.month).zfill(2)}'

def k_to_y(k: int):
    return 90 + 25 + 50*k

def add_text_task(slide: Slide, text: str, left_pt, top_pt, width_pt=0, height_pt=0, size_pt=5, rgb_color=BLUE):
    # Раскрасить
    textbox = slide.shapes.add_textbox(Pt(left_pt), Pt(top_pt), Pt(width_pt), Pt(height_pt))
    textbox.text_frame.word_wrap = True
    textbox_paragraph = textbox.text_frame.paragraphs[0] 
    textbox_paragraph.text = text
    textbox_paragraph.font.size = Pt(size_pt)
    textbox_paragraph.font.color.rgb = rgb_color
    textbox_paragraph.alignment = PP_ALIGN.LEFT
    textbox_paragraph.font.name = 'Arial'
    textbox_paragraph.font.bold = True

    return textbox_paragraph


class Report:
    def __init__(self, now_date: date) -> None:
        self.presentation = Presentation()
        self.presentation.slide_width = Pt(960)
        self.presentation.slide_height = Pt(540)
        self.slide = self.presentation.slides.add_slide( self.presentation.slide_layouts[6] )

        self.now_date = now_date
        self.year = now_date.year
        self.kvartal = (now_date.month - 1) // 3

        self.add_shape(MSO_SHAPE.RECTANGLE, 0, 60, 960, 30, color=BLUE)

        self.add_shape(MSO_SHAPE.LINE_INVERSE, 15, 60, 1, 480, color=LIGHTBLUE)
        self.add_shape(MSO_SHAPE.LINE_INVERSE, 30, 60, 1, 480, color=LIGHTBLUE)
        self.add_shape(MSO_SHAPE.LINE_INVERSE, self.date_to_x(date(1, 1, 1)), 60, 1, 480, color=LIGHTBLUE)
        self.add_shape(MSO_SHAPE.LINE_INVERSE, self.date_to_x(date(1, 4, 1)), 60, 1, 480, color=LIGHTBLUE)
        self.add_shape(MSO_SHAPE.LINE_INVERSE, self.date_to_x(date(1, 7, 1)), 60, 1, 480, color=LIGHTBLUE)
        self.add_shape(MSO_SHAPE.LINE_INVERSE, self.date_to_x(date(1, 10, 1)), 60, 1, 480, color=LIGHTBLUE)

        self.add_shape(MSO_SHAPE.LINE_INVERSE, self.date_to_x(date(1, self.kvartal * 3 + 1, 1)), 75, 540, 1, color=LIGHTBLUE)
        self.add_shape(MSO_SHAPE.LINE_INVERSE, self.date_to_x(date(1, self.kvartal * 3 + 2, 1)), 75, 1, 465, color=LIGHTBLUE)
        self.add_shape(MSO_SHAPE.LINE_INVERSE, self.date_to_x(date(1, self.kvartal * 3 + 3, 1)), 75, 1, 465, color=LIGHTBLUE)

        self.slide.shapes.add_picture('vtb_logo.png', Pt(880), Pt(10), Pt(127 // 2), Pt(52 // 2))

        self.add_text(f'Кластер «Управление Продажами», {self.kvartal + 1} суперспринт {self.year} года', 20, 0, size_pt=20, bold=True, width_pt=940, align=PP_ALIGN.LEFT)

        self.add_text('№', 0, 60, 15, 30, color=WHITE, size_pt=8, bold=True, anchor=MSO_ANCHOR.MIDDLE)
        self.add_text('К', 15, 60, 15, 30, color=WHITE, size_pt=8, bold=True, anchor=MSO_ANCHOR.MIDDLE)
        self.add_text('Задача', 30, 60, 210, 30, color=WHITE, size_pt=8, bold=True, anchor=MSO_ANCHOR.MIDDLE)
        self.add_text(f'1Q {self.year}', self.date_to_x(date(2024, 1, 1)) + 1, 60, 540 if self.kvartal == 0 else 60, 15 if self.kvartal == 0 else 30, color=WHITE, size_pt=8, bold=True, anchor=MSO_ANCHOR.MIDDLE)
        self.add_text(f'2Q {self.year}', self.date_to_x(date(2024, 4, 1)) + 1, 60, 540 if self.kvartal == 1 else 60, 15 if self.kvartal == 1 else 30, color=WHITE, size_pt=8, bold=True, anchor=MSO_ANCHOR.MIDDLE)
        self.add_text(f'3Q {self.year}', self.date_to_x(date(2024, 7, 1)) + 1, 60, 540 if self.kvartal == 2 else 60, 15 if self.kvartal == 2 else 30, color=WHITE, size_pt=8, bold=True, anchor=MSO_ANCHOR.MIDDLE)
        self.add_text(f'4Q {self.year}', self.date_to_x(date(2024, 10, 1)) + 1, 60, 540 if self.kvartal == 3 else 60, 15 if self.kvartal == 3 else 30, color=WHITE, size_pt=8, bold=True, anchor=MSO_ANCHOR.MIDDLE)

        months = ['январь', 'февраль', 'март', 'апрель', 'май', 'июнь', 'июль', 'август', 'сентябрь', 'октябрь', 'ноябрь', 'декабрь']
        for i, month in enumerate(months[3*self.kvartal:3*self.kvartal + 3]):
            self.add_text(month, self.date_to_x(date(1, 3*self.kvartal + i + 1, 1)), 75, 180, 15, color=WHITE, size_pt=6, bold=True, anchor=MSO_ANCHOR.MIDDLE)

        now_line = self.add_shape(MSO_SHAPE.LINE_INVERSE, self.date_to_x(self.now_date), 90, 1, 410, color=RED, fill=True)
        now_line.line.width = Pt(2)
        now_line.line.dash_style = MSO_LINE.DASH
        self.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, self.date_to_x(self.now_date) - 2.7, 500, 5.4, 8, color=BLUE, fill=True)
        self.add_text(date_to_text(self.now_date), self.date_to_x(self.now_date) - 20, 510, color=RED, size_pt=9, width_pt=40, height_pt=20, bold=True)

        self.add_figure(key_fig, 0, 520)
        self.add_text('ИФТ', 15, 522, width_pt=120, height_pt=16, size_pt=8, color=BLACK, align=PP_ALIGN.LEFT)
        self.add_figure(gear_fig, 40, 520)
        self.add_text('НТ', 55, 522, width_pt=120, height_pt=16, size_pt=8, color=BLACK, align=PP_ALIGN.LEFT)
        self.add_figure(gear_fig, 80, 520)
        self.add_text('ПСИ', 95, 522, width_pt=120, height_pt=16, size_pt=8, color=BLACK, align=PP_ALIGN.LEFT)
        self.add_figure(rocket_fig, 120, 520)
        self.add_text('Прод', 135, 522, width_pt=120, height_pt=16, size_pt=8, color=BLACK, align=PP_ALIGN.LEFT)
        self.add_figure(pin_fig, 160, 520)
        self.add_text('MVP', 175, 522, width_pt=120, height_pt=16, size_pt=8, color=BLACK, align=PP_ALIGN.LEFT)
        self.add_figure(pilot_fig, 200, 520)
        self.add_text('Пилот', 215, 522, width_pt=120, height_pt=16, size_pt=8, color=BLACK, align=PP_ALIGN.LEFT)

        self.add_figure(colored_flag(GRAY, filled=False), 490, 520)
        self.add_text('План', 500, 522, width_pt=200, height_pt=16, size_pt=8, color=BLACK, align=PP_ALIGN.LEFT)
        self.add_arrow(530, 525, 30, 0)
        self.add_text('Перенос срока', 565, 522, width_pt=200, height_pt=16, size_pt=8, color=BLACK, align=PP_ALIGN.LEFT)
        self.add_figure(colored_flag(GREEN, filled=True), 640, 520)
        self.add_text('Выполнено', 650, 522, width_pt=200, height_pt=16, size_pt=8, color=BLACK, align=PP_ALIGN.LEFT)
        self.add_figure(colored_flag(YELLOW, filled=False), 700, 520)
        self.add_text('Риск сдвига', 710, 522, width_pt=200, height_pt=16, size_pt=8, color=BLACK, align=PP_ALIGN.LEFT)
        self.add_figure(colored_flag(RED, filled=False), 760, 520)
        self.add_text('Просрочено', 770, 522, width_pt=200, height_pt=16, size_pt=8, color=BLACK, align=PP_ALIGN.LEFT)
        self.add_figure(colored_flag(YELLOW, filled=True), 820, 520)
        self.add_text('Выполнено со сдвигом срока', 830, 522, width_pt=200, height_pt=16, size_pt=8, color=BLACK, align=PP_ALIGN.LEFT)

    def add_shape(self, shape: MSO_SHAPE_TYPE, left_pt, top_pt, width_pt, height_pt, color=GREEN, fill=True, fill_color=None):
        shape: Shape = self.slide.shapes.add_shape(shape, Pt(left_pt), Pt(top_pt), Pt(width_pt), Pt(height_pt))
        if fill:
            shape.fill.solid()
            shape.fill.fore_color.rgb = color if fill_color is None else fill_color
        else:
            shape.fill.background()
            shape.line.width = Pt(1)

        if fill_color is not None:
            print('+')
            shape.line.width = Pt(2)
        shape.line.color.rgb = color
        return shape
    
    def add_arrow(self, left_pt, top_pt, width_pt, height_pt, color=YELLOW):
        shape = self.slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Pt(left_pt + width_pt), Pt(top_pt), Pt(left_pt), Pt(top_pt + height_pt))
        shape.line._get_or_add_ln().append(parse_xml('<a:headEnd type="arrow" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"/>'))
        shape.line.width = Pt(2)
        shape.line.dash_style = MSO_LINE.DASH
        shape.line.fill.solid()
        shape.line.color.rgb = color

        return shape
    
    def save(self, filename: str):
        self.presentation.save(filename)

    def date_to_x(self, date_: date):
        date_tuple = date_.timetuple()
        day_in_year = date_tuple.tm_yday - 1
        month_in_year = date_tuple.tm_mon - 1

        kvartal_pixels = [0] + [540 if (k == self.kvartal) else 60 for k in range(4)]
        kvartal_pixel_bounds = np.cumsum(kvartal_pixels)

        kvartal_days = [0] + [90, 91, 92, 92]
        kvartal_day_bounds = np.cumsum(kvartal_days)

        kv = month_in_year // 3

        return 240 + kvartal_pixel_bounds[kv] + (day_in_year - kvartal_day_bounds[kv]) * (kvartal_pixels[kv + 1] / kvartal_days[kv + 1])

    def add_figure(self, figure: BaseShape, left_pt, top_pt):
        element = deepcopy(figure.element)
        element.x = Pt(left_pt)
        element.y = Pt(top_pt)
        self.slide.shapes._spTree.append(element)
    
    def add_text(self, text: str, left_pt, top_pt, width_pt=100, height_pt=40, size_pt=5, color=BLUE, bold=False, anchor=MSO_ANCHOR.TOP, align=PP_ALIGN.CENTER):
        text_frame = self.slide.shapes.add_textbox(Pt(left_pt), Pt(top_pt), Pt(width_pt), Pt(height_pt)).text_frame
        text_frame.margin_left = Pt(0)
        text_frame.margin_top = Pt(0)
        text_frame.margin_right = Pt(0)
        text_frame.margin_bottom = Pt(0)
        text_frame.vertical_anchor = anchor

        textbox_paragraph = text_frame.paragraphs[0] 
        textbox_paragraph.text = text
        textbox_paragraph.font.size = Pt(size_pt)
        textbox_paragraph.font.color.rgb = color
        textbox_paragraph.alignment = align
        textbox_paragraph.font.name = 'Arial'
        textbox_paragraph.font.bold = bold

        return textbox_paragraph
    
    def add_task(self, texts: str, left_pt, top_pt, width_pt=150, height_pt=50):
        texts = [text.replace('//', '\n') for text in texts]
        text_frame = self.slide.shapes.add_textbox(Pt(left_pt), Pt(top_pt), Pt(width_pt - 30), Pt(height_pt)).text_frame
        text_frame.word_wrap = True
        text_frame.margin_left = Pt(3)
        text_frame.margin_top = Pt(1)
        text_frame.margin_right = Pt(1)
        text_frame.margin_bottom = Pt(1)
        text_frame.vertical_anchor = MSO_ANCHOR.TOP

        textbox_paragraph_0 = text_frame.paragraphs[0]
        textbox_paragraph_0.text = texts[0]
        textbox_paragraph_0.font.size = Pt(5)
        textbox_paragraph_0.font.color.rgb = BLUE
        textbox_paragraph_0.alignment = PP_ALIGN.LEFT
        textbox_paragraph_0.font.name = 'Arial'
        textbox_paragraph_0.font.bold = True

        textbox_paragraph_4 = text_frame.add_paragraph()
        textbox_paragraph_4.text = texts[1]
        textbox_paragraph_4.font.size = Pt(5)
        textbox_paragraph_4.font.color.rgb = BLUE
        textbox_paragraph_4.alignment = PP_ALIGN.LEFT
        textbox_paragraph_4.font.name = 'Arial'

        textbox_paragraph_5 = text_frame.add_paragraph()
        textbox_paragraph_5.text = texts[2]
        textbox_paragraph_5.font.size = Pt(5)
        textbox_paragraph_5.font.color.rgb = DARKGRAY
        textbox_paragraph_5.alignment = PP_ALIGN.LEFT
        textbox_paragraph_5.font.name = 'Arial'

        text_frame = self.slide.shapes.add_textbox(Pt(left_pt + width_pt - 30), Pt(top_pt), Pt(30), Pt(height_pt)).text_frame
        text_frame.word_wrap = True
        text_frame.margin_left = Pt(1)
        text_frame.margin_top = Pt(1)
        text_frame.margin_right = Pt(1)
        text_frame.margin_bottom = Pt(1)
        text_frame.vertical_anchor = MSO_ANCHOR.TOP

        textbox_paragraph_1 = text_frame.paragraphs[0]
        textbox_paragraph_1.text = texts[3]
        textbox_paragraph_1.font.size = Pt(5)
        textbox_paragraph_1.font.color.rgb = BLUE
        textbox_paragraph_1.alignment = PP_ALIGN.RIGHT
        textbox_paragraph_1.font.name = 'Arial'

        textbox_paragraph_2 = text_frame.add_paragraph()
        textbox_paragraph_2.text = texts[4]
        textbox_paragraph_2.font.size = Pt(5)
        textbox_paragraph_2.font.color.rgb = YELLOW
        textbox_paragraph_2.alignment = PP_ALIGN.RIGHT
        textbox_paragraph_2.font.name = 'Arial'

        textbox_paragraph_3 = text_frame.add_paragraph()
        textbox_paragraph_3.text = texts[5]
        textbox_paragraph_3.font.size = Pt(5)
        textbox_paragraph_3.font.color.rgb = GREEN
        textbox_paragraph_3.alignment = PP_ALIGN.RIGHT
        textbox_paragraph_3.font.name = 'Arial'

class Timeline:
    def __init__(self, report: Report, y_top: int, y_bottom: int, start_date: datetime, final_date: datetime, description: list[str] = [], show_start_date=True) -> None:
        self.report = report
        self.timeline_y = (y_top + y_bottom) / 2
        self.y_top = y_top
        self.y_bottom = y_bottom

        start_x = report.date_to_x(start_date)
        now_x = report.date_to_x(report.now_date) 
        final_x = report.date_to_x(final_date)

        if final_date < report.now_date:
            report.add_shape(MSO_SHAPE.OVAL, start_x - 6, self.timeline_y - 6, 12, 12, color=GREEN, fill=True)
            report.add_shape(MSO_SHAPE.LINE_INVERSE, start_x + 6, self.timeline_y, max(now_x - start_x - 6, 1), 0, color=GREEN)

        elif start_date < report.now_date and report.now_date <= final_date:
            report.add_shape(MSO_SHAPE.OVAL, start_x - 6, self.timeline_y - 6, 12, 12, color=GREEN, fill=True)
            report.add_shape(MSO_SHAPE.LINE_INVERSE, start_x + 6, self.timeline_y, max(now_x - start_x - 6, 1), 0, color=GREEN)
            report.add_shape(MSO_SHAPE.LINE_INVERSE, now_x, self.timeline_y, final_x - now_x, 0, color=GRAY)

        elif report.now_date <= start_date:
            report.add_shape(MSO_SHAPE.OVAL, start_x - 6, self.timeline_y - 6, 12, 12, color=GRAY, fill=True, fill_color=WHITE)
            report.add_shape(MSO_SHAPE.LINE_INVERSE, start_x + 6, self.timeline_y, max(now_x - start_x - 6, 1), 0, color=GRAY)

        report.add_shape(MSO_SHAPE.LINE_INVERSE, 30, y_bottom, 930, 0, color=LIGHTBLUE, fill=True)
        report.add_task(description, 30, y_top, 210, abs(y_bottom - y_top))

    def add_pictogram(self, figure, color: RGBColor, date_: date, write_date: bool = True, note: str = '', fill=True):
        figure = deepcopy(figure)
        if isinstance(figure, Shape):
            figure.line.color.rgb = color
            figure.fill.solid()
            if fill:
                figure.fill.fore_color.rgb = color
            else:
                figure.fill.fore_color.rgb = WHITE
        elif isinstance(figure, GroupShape):
            for subfigure in figure.shapes:
                if isinstance(subfigure, Shape):
                    subfigure.line.color.rgb = color
                    subfigure.fill.solid()
                    if fill:
                        subfigure.fill.fore_color.rgb = color
                    else:
                        subfigure.fill.fore_color.rgb = WHITE
                else:
                    try:
                        subfigure.line.color.rgb = color
                    except: 
                        pass
                    if fill:
                        try:
                            subfigure.fill.solid()
                            subfigure.fill.fore_color.rgb = color
                        except: 
                            pass
                    else:
                        try:
                            subfigure.fill.solid()
                            subfigure.fill.fore_color.rgb = WHITE
                        except: 
                            pass

            
        self.report.add_figure(figure, self.report.date_to_x(date_) - 6, self.timeline_y - 14)
        self.report.add_text(date_to_text(date_), self.report.date_to_x(date_) - 20, self.timeline_y + 3, width_pt=40, height_pt=10)
        self.report.add_text(note, self.report.date_to_x(date_) - 20, self.timeline_y - 20, width_pt=40, height_pt=10)

    def add_arrow(self, start_date: date, final_date: date):
        self.report.add_arrow(self.report.date_to_x(start_date) - 6, self.timeline_y - 3, self.report.date_to_x(final_date) - self.report.date_to_x(start_date), 0, color=YELLOW if start_date < final_date else GREEN)

    def add_gold(self, text: str, y: int):
        shape: Shape = self.report.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, 0, self.y_top + y, 30, 10, fill=True, color=BLACK, fill_color=SUNNY_YELLOW)
        shape.line.width = Pt(0)
        shape_paragraph = shape.text_frame.paragraphs[0]
        shape_paragraph.text = text
        shape_paragraph.font.size = Pt(5)
        shape_paragraph.font.name = 'Arial'
        shape_paragraph.font.color.rgb = BLACK

    def add_comment(self, text: str, date_: date, size: float):
        shape: Shape = self.report.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, self.report.date_to_x(date_) - size / 2, self.timeline_y + 10, size, 10, fill=False, color=YELLOW)
        shape_paragraph = shape.text_frame.paragraphs[0]
        shape_paragraph.text = text
        shape_paragraph.font.size = Pt(5)
        shape_paragraph.font.name = 'Arial'
        shape_paragraph.font.color.rgb = BLUE

