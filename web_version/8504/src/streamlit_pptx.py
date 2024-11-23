import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE, MSO_CONNECTOR_TYPE, MSO_CONNECTOR
from pptx.dml.color import RGBColor
from pptx.slide import SlideMasters, Slides, Slide, SlideLayout
from pptx.util import Inches, Pt, Cm, Emu
from copy import deepcopy
from pptx.shapes.group import GroupShape
from pptx.shapes.base import BaseShape
from pptx.shapes.connector import Connector
from pptx.shapes.freeform import FreeformBuilder
from pptx.shapes.autoshape import Shape
from collections import defaultdict
from lxml import etree
import numpy as np
from datetime import datetime, date, timedelta
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.shapes.autoshape import AutoShapeType, Shape, BaseShape
from pptx.enum.dml import MSO_LINE
from pptx.shapes.graphfrm import GraphicFrame
from pptx.text.text import _Paragraph
from pptx.oxml import parse_xml


CONTENT_DIR = '../contents/'


slide_template_path = os.path.join(CONTENT_DIR, 'examples.pptx')
slide_template = Presentation(slide_template_path).slides[0]

# Загружаем необходимые фигуры из шаблона
base_flag = deepcopy(slide_template.shapes[0])
green_flag = deepcopy(slide_template.shapes[0])
risky_flag = deepcopy(slide_template.shapes[1])
failed_flag = deepcopy(slide_template.shapes[3])
finally_flag = deepcopy(slide_template.shapes[8])
key_fig = deepcopy(slide_template.shapes[10])
gear_fig = deepcopy(slide_template.shapes[11])
box_fig = deepcopy(slide_template.shapes[12])
rocket_fig = deepcopy(slide_template.shapes[13])
pin_fig = deepcopy(slide_template.shapes[14])
pilot_fig = deepcopy(slide_template.shapes[15])
plan_flag = deepcopy(slide_template.shapes[16])

# Цвета
GREEN = RGBColor(78, 201, 149)
BLUE = RGBColor(10, 40, 150)
LIGHTBLUE = RGBColor(200, 200, 255)
YELLOW = RGBColor(240, 160, 40)
RED = RGBColor(255, 40, 40)
GRAY = RGBColor(194, 194, 194)
WHITE = RGBColor(255, 255, 255)
BLACK = RGBColor(0, 0, 0)
DARKGRAY = RGBColor(106, 106, 106)
SUNNY_YELLOW = RGBColor(255, 255, 0)
FOUNDING = RGBColor(125, 150, 246)
COMPLEXITY = RGBColor(67, 104, 172)

base_flag.width = Pt(8.4)
base_flag.height = Pt(12.2)

def colored_flag(color: RGBColor, filled: bool):
    flag: GroupShape = deepcopy(base_flag)
    fabric: Shape = flag.shapes[0]
    stick: Shape = flag.shapes[1]

    fabric.line.color.rgb = color
    fabric.fill.solid()
    fabric.fill.fore_color.rgb = color if filled else WHITE
    stick.line.color.rgb = color

    return flag

def date_to_text(date_: date):
    return f'{str(date_.day).zfill(2)}.{str(date_.month).zfill(2)}'

class Report:
    def __init__(self, now_date: date, quarters_with_tasks, tasks_per_quarter, total_tasks, min_quarter_width=60, max_quarter_width=300, left_margin=240, right_margin=20) -> None:
        self.presentation = Presentation()
        self.presentation.slide_width = Pt(960)
        self.presentation.slide_height = Pt(540)
        self.slide = self.presentation.slides.add_slide(self.presentation.slide_layouts[6])

        self.now_date = now_date
        self.quarters_with_tasks = quarters_with_tasks
        self.tasks_per_quarter = tasks_per_quarter
        self.total_tasks = total_tasks

        # Создаём список всех кварталов с задачами
        self.quarter_list = quarters_with_tasks

        # Доступная ширина для таймлайна
        available_width = 960 - left_margin - right_margin

        # Рассчитываем суммарный вес кварталов на основе количества задач
        total_weight = sum(tasks_per_quarter.values())

        # Рассчитываем ширину каждого квартала пропорционально его весу (количеству задач)
        self.kvartal_pixels = []
        for quarter in self.quarter_list:
            weight = tasks_per_quarter[quarter] / total_weight
            width = available_width * weight
            width = max(min_quarter_width, min(width, max_quarter_width))
            self.kvartal_pixels.append(width)

        # Если суммарная ширина кварталов превышает доступную ширину, масштабируем их
        total_quarter_width = sum(self.kvartal_pixels)
        if total_quarter_width > available_width:
            scaling_factor = available_width / total_quarter_width
            self.kvartal_pixels = [width * scaling_factor for width in self.kvartal_pixels]
        else:
            # Если осталось свободное пространство, распределяем его равномерно
            extra_space = available_width - total_quarter_width
            self.kvartal_pixels = [width + extra_space / len(self.kvartal_pixels) for width in self.kvartal_pixels]

        self.kvartal_pixel_bounds = np.cumsum([0] + self.kvartal_pixels)

        # Расчет дней в каждом квартале
        self.kvartal_days = []
        for year, quarter in self.quarter_list:
            q_start_month = (quarter - 1) * 3 + 1
            q_end_month = q_start_month + 2
            q_start = date(year, q_start_month, 1)
            if q_end_month >= 12:
                q_end = date(year, 12, 31)
            else:
                q_end = date(year, q_end_month + 1, 1) - timedelta(days=1)
            days_in_quarter = (q_end - q_start).days + 1
            self.kvartal_days.append(days_in_quarter)
        self.kvartal_day_bounds = np.cumsum([0] + self.kvartal_days)

        # Сохраняем начальный год и квартал
        self.start_year = self.quarter_list[0][0]
        self.start_quarter = self.quarter_list[0][1]

        # Устанавливаем левый отступ
        self.left_margin = left_margin

        # Создаем основную структуру отчета
        self.create_base_structure()

    def create_base_structure(self):
        self.add_shape(MSO_SHAPE.RECTANGLE, 0, 60, 960, 30, color=BLUE)

        # Рисуем вертикальные линии для каждого квартала
        for i in range(len(self.kvartal_pixel_bounds)):
            x_position = self.left_margin + self.kvartal_pixel_bounds[i]
            self.add_shape(MSO_SHAPE.LINE_INVERSE, x_position, 60, 1, 480, color=LIGHTBLUE)

        # Рисуем горизонтальные линии
        self.add_shape(MSO_SHAPE.LINE_INVERSE, 15, 60, 1, 480, color=LIGHTBLUE)
        self.add_shape(MSO_SHAPE.LINE_INVERSE, 30, 60, 1, 480, color=LIGHTBLUE)

        # Добавляем логотип с корректным путем
        logo_path = os.path.join(CONTENT_DIR, 'vtb_logo.png')
        self.slide.shapes.add_picture(logo_path, Pt(880), Pt(10), Pt(127 // 2), Pt(52 // 2))

        self.add_text(f'Кластер «Управление Продажами»', 20, 0, size_pt=20, bold=True, width_pt=940, align=PP_ALIGN.LEFT)

        self.add_text('№', 0, 60, 15, 30, color=WHITE, size_pt=8, bold=True, anchor=MSO_ANCHOR.MIDDLE)
        self.add_text('К', 15, 60, 15, 30, color=WHITE, size_pt=8, bold=True, anchor=MSO_ANCHOR.MIDDLE)
        self.add_text('Задача', 30, 60, 210, 30, color=WHITE, size_pt=8, bold=True, anchor=MSO_ANCHOR.MIDDLE)

        # Добавляем названия кварталов
        for i, (year, quarter) in enumerate(self.quarter_list):
            quarter_name = f'{quarter}Q {year}'
            x_position = self.left_margin + self.kvartal_pixel_bounds[i]
            width = self.kvartal_pixels[i]
            self.add_text(quarter_name, x_position + 1, 60, width, 15, color=WHITE, size_pt=8, bold=True, anchor=MSO_ANCHOR.MIDDLE)

        # Добавляем названия месяцев
        months = ['январь', 'февраль', 'март', 'апрель', 'май', 'июнь', 'июль', 'август', 'сентябрь', 'октябрь', 'ноябрь', 'декабрь']
        for i, (year, quarter) in enumerate(self.quarter_list):
            for m in range(3):
                month_index = (quarter - 1) * 3 + m
                if month_index >= 12:
                    break
                month_name = months[month_index]
                month_start_date = date(year, month_index + 1, 1)
                x_position = self.date_to_x(month_start_date)
                self.add_text(month_name, x_position, 75, self.kvartal_pixels[i]/3, 15, color=WHITE, size_pt=6, bold=True, anchor=MSO_ANCHOR.MIDDLE)

        # Рисуем текущую дату, если она в диапазоне
        try:
            now_x = self.date_to_x(self.now_date)
            now_line = self.add_shape(MSO_SHAPE.LINE_INVERSE, now_x, 90, 1, 410, color=RED, fill=True)
            now_line.line.width = Pt(2)
            now_line.line.dash_style = MSO_LINE.DASH
            self.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, now_x - 2.7, 500, 5.4, 8, color=BLUE, fill=True)
            self.add_text(date_to_text(self.now_date), now_x - 20, 510, color=RED, size_pt=9, width_pt=40, height_pt=20, bold=True)
        except ValueError:
            pass  # Если текущая дата вне диапазона, не отображаем линию

        # Добавляем легенду
        self.add_legend()

    def add_legend(self):
        # Добавляем фигуры и пояснения для легенды
        self.add_figure(key_fig, 0, 520)
        self.add_text('ИФТ', 15, 522, width_pt=120, height_pt=16, size_pt=8, color=BLACK, align=PP_ALIGN.LEFT)
        self.add_figure(gear_fig, 40, 520)
        self.add_text('НТ', 55, 522, width_pt=120, height_pt=16, size_pt=8, color=BLACK, align=PP_ALIGN.LEFT)
        self.add_figure(box_fig, 80, 520)
        self.add_text('ПСИ', 95, 522, width_pt=120, height_pt=16, size_pt=8, color=BLACK, align=PP_ALIGN.LEFT)
        self.add_figure(rocket_fig, 120, 520)
        self.add_text('Прод', 135, 522, width_pt=120, height_pt=16, size_pt=8, color=BLACK, align=PP_ALIGN.LEFT)
        self.add_figure(pin_fig, 160, 520)
        self.add_text('MVP', 175, 522, width_pt=120, height_pt=16, size_pt=8, color=BLACK, align=PP_ALIGN.LEFT)
        self.add_figure(pilot_fig, 200, 520)
        self.add_text('Пилот', 215, 522, width_pt=120, height_pt=16, size_pt=8, color=BLACK, align=PP_ALIGN.LEFT)

        self.add_figure(colored_flag(GRAY, filled=False), 490, 520)
        self.add_text('План', 500, 522, width_pt=200, height_pt=16, size_pt=8, color=BLACK, align=PP_ALIGN.LEFT)
        self.add_arrow(530, 525, 30, 0)
        self.add_text('Перенос срока', 565, 522, width_pt=200, height_pt=16, size_pt=8, color=BLACK, align=PP_ALIGN.LEFT)
        self.add_figure(colored_flag(GREEN, filled=True), 640, 520)
        self.add_text('Выполнено', 650, 522, width_pt=200, height_pt=16, size_pt=8, color=BLACK, align=PP_ALIGN.LEFT)
        self.add_figure(colored_flag(YELLOW, filled=False), 700, 520)
        self.add_text('Риск сдвига', 710, 522, width_pt=200, height_pt=16, size_pt=8, color=BLACK, align=PP_ALIGN.LEFT)
        self.add_figure(colored_flag(RED, filled=False), 760, 520)
        self.add_text('Просрочено', 770, 522, width_pt=200, height_pt=16, size_pt=8, color=BLACK, align=PP_ALIGN.LEFT)
        self.add_figure(colored_flag(YELLOW, filled=True), 820, 520)
        self.add_text('Выполнено со сдвигом срока', 830, 522, width_pt=200, height_pt=16, size_pt=8, color=BLACK, align=PP_ALIGN.LEFT)

    def add_shape(self, shape: MSO_SHAPE_TYPE, left_pt, top_pt, width_pt, height_pt, color=GREEN, fill=True, fill_color=None):
        shape_obj = self.slide.shapes.add_shape(shape, Pt(left_pt), Pt(top_pt), Pt(width_pt), Pt(height_pt))
        if fill:
            shape_obj.fill.solid()
            shape_obj.fill.fore_color.rgb = color if fill_color is None else fill_color
        else:
            shape_obj.fill.background()
            shape_obj.line.width = Pt(1)

        if fill_color is not None:
            shape_obj.line.width = Pt(2)
        shape_obj.line.color.rgb = color
        return shape_obj

    def add_arrow(self, left_pt, top_pt, width_pt, height_pt, color=YELLOW):
        shape = self.slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Pt(left_pt + width_pt), Pt(top_pt), Pt(left_pt), Pt(top_pt + height_pt))
        shape.line._get_or_add_ln().append(parse_xml('<a:headEnd type="arrow" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"/>'))
        shape.line.width = Pt(2)
        shape.line.dash_style = MSO_LINE.DASH
        shape.line.fill.solid()
        shape.line.color.rgb = color

        return shape

    def save(self, filename: str):
        output_path = os.path.join(CONTENT_DIR, filename)
        self.presentation.save(output_path)

    def date_to_x(self, date_: date):
        # Определяем индекс квартала
        try:
            kv = self.quarter_list.index((date_.year, (date_.month - 1) // 3 + 1))
        except ValueError:
            raise ValueError("Date out of range")

        kvartal_pixel = self.kvartal_pixels[kv]
        kvartal_day = self.kvartal_days[kv]

        q_year, q_quarter = self.quarter_list[kv]
        q_start_month = (q_quarter - 1) * 3 + 1
        q_start_date = date(q_year, q_start_month, 1)
        days_into_quarter = (date_ - q_start_date).days

        x = self.left_margin + self.kvartal_pixel_bounds[kv] + days_into_quarter * (kvartal_pixel / kvartal_day)
        return x

    def add_figure(self, figure: BaseShape, left_pt, top_pt):
        element = deepcopy(figure.element)
        element.x = Pt(left_pt)
        element.y = Pt(top_pt)
        self.slide.shapes._spTree.append(element)

    def add_text(self, text: str, left_pt, top_pt, width_pt=100, height_pt=40, size_pt=5, color=BLUE, bold=False, anchor=MSO_ANCHOR.TOP, align=PP_ALIGN.CENTER):
        text_frame = self.slide.shapes.add_textbox(Pt(left_pt), Pt(top_pt), Pt(width_pt), Pt(height_pt)).text_frame
        text_frame.margin_left = Pt(0)
        text_frame.margin_top = Pt(0)
        text_frame.margin_right = Pt(0)
        text_frame.margin_bottom = Pt(0)
        text_frame.vertical_anchor = anchor

        textbox_paragraph = text_frame.paragraphs[0]
        textbox_paragraph.text = text
        textbox_paragraph.font.size = Pt(size_pt)
        textbox_paragraph.font.color.rgb = color
        textbox_paragraph.alignment = align
        textbox_paragraph.font.name = 'Arial'
        textbox_paragraph.font.bold = bold

        return textbox_paragraph

    def add_task(self, texts: list[str], left_pt, top_pt, width_pt=150, height_pt=50):
        texts = [text.replace('//', '\n') for text in texts]
        text_frame = self.slide.shapes.add_textbox(Pt(left_pt), Pt(top_pt), Pt(width_pt - 30), Pt(height_pt)).text_frame
        text_frame.word_wrap = True
        text_frame.margin_left = Pt(3)
        text_frame.margin_top = Pt(1)
        text_frame.margin_right = Pt(1)
        text_frame.margin_bottom = Pt(1)
        text_frame.vertical_anchor = MSO_ANCHOR.TOP

        for idx, text in enumerate(texts[:3]):
            paragraph = text_frame.paragraphs[idx] if idx == 0 else text_frame.add_paragraph()
            paragraph.text = text
            paragraph.font.size = Pt(5)
            paragraph.font.color.rgb = [BLUE, BLUE, DARKGRAY][idx]
            paragraph.alignment = PP_ALIGN.LEFT
            paragraph.font.name = 'Arial'
            paragraph.font.bold = idx == 0

        text_frame_right = self.slide.shapes.add_textbox(Pt(left_pt + width_pt - 30), Pt(top_pt), Pt(30), Pt(height_pt)).text_frame
        text_frame_right.word_wrap = True
        text_frame_right.margin_left = Pt(1)
        text_frame_right.margin_top = Pt(1)
        text_frame_right.margin_right = Pt(1)
        text_frame_right.margin_bottom = Pt(1)
        text_frame_right.vertical_anchor = MSO_ANCHOR.TOP

        for idx, text in enumerate(texts[3:]):
            paragraph = text_frame_right.paragraphs[idx] if idx == 0 else text_frame_right.add_paragraph()
            paragraph.text = text
            paragraph.font.size = Pt(5)
            paragraph.font.color.rgb = [BLUE, YELLOW, GREEN][idx]
            paragraph.alignment = PP_ALIGN.RIGHT
            paragraph.font.name = 'Arial'

class Timeline:
    def __init__(self, report: Report, y_top: int, y_bottom: int, start_date: date, final_date: date, description: list[str] = [], show_start_date=True) -> None:
        self.report = report
        self.timeline_y = (y_top + y_bottom) / 2
        self.y_top = y_top
        self.y_bottom = y_bottom

        try:
            start_x = report.date_to_x(start_date)
            final_x = report.date_to_x(final_date)
        except ValueError:
            # Если даты вне диапазона, устанавливаем их в границы таймлайна
            if start_date < report.now_date:
                start_x = report.left_margin
            else:
                start_x = report.left_margin + sum(report.kvartal_pixels)
            if final_date > report.now_date:
                final_x = report.left_margin + sum(report.kvartal_pixels)
            else:
                final_x = report.left_margin

        # Проверяем, находится ли текущая дата в диапазоне
        try:
            now_x = report.date_to_x(report.now_date)
        except ValueError:
            now_x = None

        if final_date < report.now_date:
            report.add_shape(MSO_SHAPE.OVAL, start_x - 6, self.timeline_y - 6, 12, 12, color=GREEN, fill=True)
            report.add_shape(MSO_SHAPE.LINE_INVERSE, start_x + 6, self.timeline_y, final_x - start_x - 6, 0, color=GREEN)
        elif start_date < report.now_date <= final_date:
            report.add_shape(MSO_SHAPE.OVAL, start_x - 6, self.timeline_y - 6, 12, 12, color=GREEN, fill=True)
            if now_x:
                report.add_shape(MSO_SHAPE.LINE_INVERSE, start_x + 6, self.timeline_y, now_x - start_x - 6, 0, color=GREEN)
                report.add_shape(MSO_SHAPE.LINE_INVERSE, now_x, self.timeline_y, final_x - now_x, 0, color=GRAY)
            else:
                report.add_shape(MSO_SHAPE.LINE_INVERSE, start_x + 6, self.timeline_y, final_x - start_x - 6, 0, color=GRAY)
        elif report.now_date <= start_date:
            report.add_shape(MSO_SHAPE.OVAL, start_x - 6, self.timeline_y - 6, 12, 12, color=GRAY, fill=True, fill_color=WHITE)
            report.add_shape(MSO_SHAPE.LINE_INVERSE, start_x + 6, self.timeline_y, final_x - start_x - 6, 0, color=GRAY)

        if show_start_date:
            report.add_text(date_to_text(start_date), start_x - 16, self.timeline_y + 10, 32, 10)
        report.add_text(date_to_text(final_date), final_x - 16, self.timeline_y + 10, 32, 10)

        report.add_shape(MSO_SHAPE.LINE_INVERSE, 30, y_top, 930, 0, color=LIGHTBLUE, fill=True)
        report.add_task(description, 30, y_top, 210, abs(y_bottom - y_top))

    def add_pictogram(self, figure, color: RGBColor, date_: date, write_date: bool = True, note: str = '', fill=True):
        figure = deepcopy(figure)
        if isinstance(figure, Shape):
            figure.line.color.rgb = color
            figure.fill.solid()
            if fill:
                figure.fill.fore_color.rgb = color
            else:
                figure.fill.fore_color.rgb = WHITE
        elif isinstance(figure, GroupShape):
            for subfigure in figure.shapes:
                if isinstance(subfigure, Shape):
                    subfigure.line.color.rgb = color
                    subfigure.fill.solid()
                    if fill:
                        subfigure.fill.fore_color.rgb = color
                    else:
                        subfigure.fill.fore_color.rgb = WHITE

        try:
            x_position = self.report.date_to_x(date_)
        except ValueError:
            # Если дата вне диапазона, пропускаем добавление пиктограммы
            return

        self.report.add_figure(figure, x_position - 6, self.timeline_y - 14)
        if write_date:
            self.report.add_text(date_to_text(date_), x_position - 20, self.timeline_y + 3, width_pt=40, height_pt=10)
        if note:
            self.report.add_text(note, x_position - 20, self.timeline_y - 20, width_pt=40, height_pt=10)

    def add_arrow(self, start_date: date, final_date: date):
        try:
            start_x = self.report.date_to_x(start_date)
            final_x = self.report.date_to_x(final_date)
        except ValueError:
            return  # Если даты вне диапазона, пропускаем добавление стрелки
        self.report.add_arrow(start_x - 6, self.timeline_y - 3, final_x - start_x, 0, color=YELLOW if start_date < final_date else GREEN)

    def add_gold(self, text: str, y: int):
        shape = self.report.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, 0, self.y_top + y, 30, 10, fill=True, color=BLACK, fill_color=SUNNY_YELLOW)
        shape.line.width = Pt(0)
        shape_paragraph = shape.text_frame.paragraphs[0]
        shape_paragraph.text = text
        shape_paragraph.font.size = Pt(5)
        shape_paragraph.font.name = 'Arial'
        shape_paragraph.font.color.rgb = BLACK

    def add_comment(self, text: str, date_: date, size: float):
        try:
            x_position = self.report.date_to_x(date_)
        except ValueError:
            return  # Если дата вне диапазона, пропускаем добавление комментария
        shape = self.report.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x_position - size / 2, self.timeline_y + 10, size, 10, fill=False, color=YELLOW)
        shape_paragraph = shape.text_frame.paragraphs[0]
        shape_paragraph.text = text
        shape_paragraph.font.size = Pt(5)
        shape_paragraph.font.name = 'Arial'
        shape_paragraph.font.color.rgb = BLUE